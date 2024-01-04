import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from tabulate import tabulate  # Install using: pip install tabulate
import tabula

import shutil

def read_content_from_file(file_path):
    # Check if the file path exists
    if not os.path.exists(file_path):
        print(f"File '{file_path}' does not exist.")
        return

    # Create a dictionary to store content from the file
    content_info = {}

    try:
        if file_path.endswith('.pdf'):
            text, tables, images = read_content_from_pdf(file_path)
        elif file_path.endswith('.pptx'):
            text, tables = read_content_from_pptx(file_path)
        elif file_path.endswith('.docx'):
            text, tables = read_content_from_docx(file_path)
        else:
            print(f"Unsupported file type: '{file_path}'")
            return

        # Store the text, tables, and images in the dictionary
        content_info['text'] = text
        content_info['tables'] = tables
        content_info['images'] = images
    except Exception as e:
        print(f"Error reading '{file_path}': {e}\n")

    return content_info

def read_content_from_pdf(file_path):
    text = ""
    tables = []
    images = []

    

    with fitz.open(file_path) as pdf_document:
        save_path = "./Images/"+f"{file_path}"
        if os.path.exists(save_path):
            # If it exists, delete the folder
            shutil.rmtree(save_path)
        
        os.makedirs(save_path)

        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            text += page.get_text()

            

            # Extract images from the page
            for img_index, img in enumerate(page.get_images(full=True)):
                img_index += 1
                img_index_str = f"{page_number + 1}_{img_index}"
                img_file_path = save_path+ f"/image_{img_index_str}.png"
                img_data = pdf_document.extract_image(img[0])
                with open(img_file_path, "wb") as img_file:
                    img_file.write(img_data["image"])
                images.append(img_file_path)



            

    return text, tables, images

def read_content_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""
    tables = []  # Tables not implemented for PPTX in this example
    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text, tables

def read_content_from_docx(file_path):
    doc = Document(file_path)
    paragraphs = [paragraph.text for paragraph in doc.paragraphs]

    tables = []
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        tables.append(table_data)

    return '\n'.join(paragraphs), tables

# Replace 'your_file_path' with the actual path of the file
def return_text(file_path):
    content_info = read_content_from_file(file_path)
    text = ""

    # Access text, tables, and images for the specified file
    if content_info:
        text += f"Contents of '{file_path}':\n{content_info['text']}\n"

        for table_number, table_data in enumerate(content_info['tables']):
            text += f"Table {table_number + 1}:\n{tabulate(table_data, headers='firstrow', tablefmt='grid')}\n"

        text += "Images:\n"
        for img_path in content_info['images']:
            text += f"{img_path}\n"

        text += "================================================================================================\n"

    return text, content_info['images']


