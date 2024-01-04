# app.py
from flask import Flask, render_template, request
import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from tabulate import tabulate  # Install using: pip install tabulate

import shutil

def read_content_from_file(file):
    # Check if the file object is valid
    if not file:
        print("Invalid file object.")
        return

    # Create a dictionary to store content from the file
    content_info = {}

    try:
        file_extension = os.path.splitext(file.name)[1].lower()

        if file_extension == '.pdf':
            text, tables, images = read_content_from_pdf(file)
        elif file_extension == '.pptx':
            text, tables = read_content_from_pptx(file)
        elif file_extension == '.docx':
            text, tables = read_content_from_docx(file)
        else:
            print(f"Unsupported file type: '{file_extension}'")
            return

        # Store the text, tables, and images in the dictionary
        content_info['text'] = text
        content_info['tables'] = tables
        content_info['images'] = images
    except Exception as e:
        print(f"Error reading '{file.name}': {e}\n")

    return content_info


def read_content_from_pdf(file):
    text = ""
    tables = []  # Tables not implemented for PDF in this example
    images = []

    with fitz.open(file.name) as pdf_document:
        save_path = "./Images/"+f"{file.name}"
        if os.path.exists(save_path):
            # If it exists, delete the folder
            shutil.rmtree(save_path)
        
        os.makedirs(save_path)

        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            text += page.get_text()

            # Extract images from the page
            for img_index, img in enumerate(page.get_images(full=True)):
                img_index += 1
                img_index_str = f"{page_number + 1}_{img_index}"
                img_file_path = save_path + f"/image_{img_index_str}.png"
                img_data = pdf_document.extract_image(img[0])
                with open(img_file_path, "wb") as img_file:
                    img_file.write(img_data["image"])
                images.append(img_file_path)

    return text, tables, images

def read_content_from_pptx(file):
    presentation = Presentation(file)
    text = ""
    tables = []  # Tables not implemented for PPTX in this example

    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text, tables

def read_content_from_docx(file):
    doc = Document(file)
    paragraphs = [paragraph.text for paragraph in doc.paragraphs]

    tables = []
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        tables.append(table_data)

    return '\n'.join(paragraphs), tables

# Replace 'your_file_path' with the actual path of the file
def return_text(file):
    content_info = read_content_from_file(file)
    text = ""

    # Access text, tables, and images for the specified file
    if content_info:
        text += f"Contents of '{file.name}':\n{content_info['text']}\n"

        for table_number, table_data in enumerate(content_info['tables']):
            text += f"Table {table_number + 1}:\n{tabulate(table_data, headers='firstrow', tablefmt='grid')}\n"

        text += "Images:\n"
        for img_path in content_info['images']:
            text += f"{img_path}\n"

        text += "================================================================================================\n"

    return text, content_info['images']


from to_text_converter import return_text
from file_to_text_converter import return_text_for_single_file

import textwrap

import google.generativeai as genai

from IPython.display import display
from IPython.display import Markdown
import PIL.Image


def to_markdown(text):
  text = text.replace('•', '  *')
  return textwrap.indent(text, '> ', predicate=lambda _: True)

GOOGLE_API_KEY= "AIzaSyDDSwzN5o85ckkRVJXZEidq9zIPKIP8HtY"

genai.configure(api_key=GOOGLE_API_KEY)

model_text = genai.GenerativeModel('gemini-pro')
model_image = genai.GenerativeModel('gemini-pro-vision')

def get_questions(file):
  answer_text = ""

  text , images = return_text(file)
  
  for img_path in images:
    img = PIL.Image.open(img_path)

    response = model_image.generate_content(img)

    answer_text += to_markdown(response.text)
  
  response1 = model_text.generate_content(text +"\n"  + "\n\n" + "make questions and answers as many as you can from above text and additional things ralated to above with answers")

  output = to_markdown(response1.text)

  output += "\n\n"

  for img_path in images:
    img = PIL.Image.open(img_path)

    response2 = model_image.generate_content(["make questions and answers as many as you can from the image and additional things ralated to above with answers", img], stream=True)
    response2.resolve()
    output += to_markdown(response2.text)
    


  return output


app = Flask(__name__)


@app.route('/', methods=['GET', 'POST'])
def index():
    result = None
    if request.method == 'POST':
        if 'file' in request.files:
            input_file = request.files['file']
            result = get_questions(input_file)
    return render_template('index.html', result=result)

if __name__ == '__main__':
    app.run(debug=True)
